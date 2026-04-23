from io import BytesIO
from pathlib import Path

from pptx import Presentation

from models import BvaCalculations, BvaRequest
from slide_updaters import (
    update_slide_01,
    update_slide_05,
    update_slide_07,
    update_slide_08,
    update_slide_09,
    update_slide_12,
    update_slide_13,
    update_slide_14,
    update_slide_15,
    update_slide_16,
)

# Docker: template is alongside Python files; local dev: template is in project root
_HERE = Path(__file__).resolve().parent
TEMPLATE_PATH = _HERE / "Db2GeniusHub_Proposal.pptx"
if not TEMPLATE_PATH.exists():
    TEMPLATE_PATH = _HERE.parent / "Db2GeniusHub_Proposal.pptx"


def generate_pptx(inputs: BvaRequest, calcs: BvaCalculations) -> bytes:
    prs = Presentation(str(TEMPLATE_PATH))

    update_slide_01(prs.slides[0], inputs)
    update_slide_05(prs.slides[4], inputs)
    update_slide_07(prs.slides[6], inputs, calcs)
    update_slide_08(prs.slides[7], inputs, calcs)
    update_slide_09(prs.slides[8], inputs, calcs)
    update_slide_12(prs.slides[11], inputs)
    update_slide_13(prs.slides[12], inputs, calcs)
    update_slide_14(prs.slides[13], inputs, calcs)
    update_slide_15(prs.slides[14], inputs, calcs)
    update_slide_16(prs.slides[15], inputs, calcs)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
