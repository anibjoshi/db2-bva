from lxml import etree
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from calculator import RAMP_YR1, RAMP_YR2, RAMP_YR3
from formatting import (
    fmt_currency_full,
    fmt_currency_short,
    fmt_percentage,
    fmt_roi,
)
from models import BvaCalculations, BvaRequest
from trade_up_catalog import calculate_line_item


def _set_cell(table, row, col, text):
    """Update a table cell's first run text, preserving formatting."""
    cell = table.cell(row, col)
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            run.text = str(text)
            return
    cell.text_frame.paragraphs[0].text = str(text)


def _set_source_cell(table, row, col, text, *, header=False):
    """Write a source citation with 24pt centered text (white if header row)."""
    cell = table.cell(row, col)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = cell.text_frame.paragraphs[0]
    para.text = str(text)
    para.alignment = PP_ALIGN.CENTER
    color = RGBColor(0xFF, 0xFF, 0xFF) if header else RGBColor(0x39, 0x39, 0x39)
    for run in para.runs:
        run.font.size = Pt(24)
        run.font.color.rgb = color
        run.font.name = "IBM Plex Sans"


# ---------------------------------------------------------------------------
# Slide 1 (index 0): Title — customer name + date
# ---------------------------------------------------------------------------
def update_slide_01(slide, inputs: BvaRequest):
    subtitle = slide.shapes[2].text_frame.paragraphs[1].runs[0]
    subtitle.text = subtitle.text.replace("Cigna", inputs.customer_name)

    date_box = slide.shapes[3]
    date_box.text_frame.paragraphs[0].runs[0].text = inputs.report_date


# ---------------------------------------------------------------------------
# Slide 5 (index 4): Trade-up Pricing
# ---------------------------------------------------------------------------
# Template layout:
#   shape[1] Title ("Cigna Trade-up Pricing")
#   shape[2] Deal-level table: [Current S&S | Support Due | Notes], 1 data row
#   shape[3] Line-item table:  [Sr. | Existing entitlement | AI Edition | Qty | Price], 3 data rows
# We replace placeholder cell text with real values and adjust line-item row count.


def _clone_row(tbl_el, src_row_idx: int) -> None:
    from copy import deepcopy
    src_tr = tbl_el.findall(qn("a:tr"))[src_row_idx]
    tbl_el.append(deepcopy(src_tr))


def _delete_row(tbl_el, row_idx: int) -> None:
    tr = tbl_el.findall(qn("a:tr"))[row_idx]
    tbl_el.remove(tr)


def _set_cell_lines(cell, lines: tuple[str, ...]) -> None:
    """Write each string in `lines` as its own paragraph in the cell.

    Preserves first-run formatting of each existing paragraph, removes extra
    runs within a paragraph, and drops any paragraphs beyond len(lines).
    """
    tf = cell.text_frame
    tf_el = tf._txBody
    existing_paras = list(tf.paragraphs)

    for i, line in enumerate(lines):
        if i < len(existing_paras):
            para = existing_paras[i]
            if para.runs:
                para.runs[0].text = str(line)
                p_el = para._p
                for extra in p_el.findall(qn("a:r"))[1:]:
                    p_el.remove(extra)
            else:
                para.text = str(line)
        else:
            # More lines than existing paragraphs — clone the last existing para
            from copy import deepcopy
            new_p = deepcopy(existing_paras[-1]._p)
            tf_el.append(new_p)
            # update its first run's text
            runs = new_p.findall(qn("a:r"))
            if runs:
                t_el = runs[0].find(qn("a:t"))
                if t_el is not None:
                    t_el.text = str(line)

    # Drop any extra paragraphs beyond len(lines)
    for p_el in tf_el.findall(qn("a:p"))[len(lines):]:
        tf_el.remove(p_el)


def _set_cell_text(cell, text: str) -> None:
    _set_cell_lines(cell, (text,))


def update_slide_05(slide, inputs: BvaRequest):
    # Title: "Cigna Trade-up Pricing" → "<Customer> Trade-up Pricing"
    title_runs = slide.shapes[1].text_frame.paragraphs[0].runs
    title_runs[0].text = title_runs[0].text.replace("Cigna", inputs.customer_name)

    # No line items: leave placeholder tables as-is.
    if not inputs.trade_up_items:
        return

    # --- Deal-level table (shape 2): Current S&S | Support Due | Notes ---
    deal_table = slide.shapes[2].table
    _set_cell_text(deal_table.cell(1, 0), fmt_currency_full(inputs.current_s_and_s_total))
    _set_cell_text(deal_table.cell(1, 1), inputs.renewal_date.strip())

    # Notes: split on newlines so each line becomes its own paragraph.
    note_lines = tuple(line.strip() for line in inputs.trade_up_notes.splitlines() if line.strip())
    _set_cell_lines(deal_table.cell(1, 2), note_lines or ("",))

    # --- Line-item table (shape 3): Sr. | Existing | AI Edition | Qty | Price ---
    li_table = slide.shapes[3].table
    li_tbl_el = li_table._tbl

    n_items = len(inputs.trade_up_items)
    existing_data_rows = len(li_table.rows) - 1  # minus header

    if n_items > existing_data_rows:
        for _ in range(n_items - existing_data_rows):
            _clone_row(li_tbl_el, existing_data_rows)  # clone last data row
    elif n_items < existing_data_rows:
        for _ in range(existing_data_rows - n_items):
            _delete_row(li_tbl_el, len(li_table.rows) - 1)  # drop from bottom

    for r, item in enumerate(inputs.trade_up_items, start=1):
        calc = calculate_line_item(item.source_product, item.source_quantity, item.discount_pct)
        _set_cell_text(li_table.cell(r, 0), str(r))
        _set_cell_text(li_table.cell(r, 1), calc.source)
        _set_cell_lines(li_table.cell(r, 2), (calc.target, calc.pn))
        _set_cell_text(li_table.cell(r, 3), str(item.source_quantity))
        _set_cell_text(li_table.cell(r, 4), fmt_currency_full(calc.year1_total))


# ---------------------------------------------------------------------------
# Slide 7 (index 6): Executive Summary (Value Proposition)
# ---------------------------------------------------------------------------
def update_slide_07(slide, inputs: BvaRequest, calcs: BvaCalculations):
    shapes = slide.shapes

    # Value Proposition intro — "<Customer> can realize the following value:"
    intro = shapes[0].text_frame.paragraphs[1].runs[0]
    intro.text = intro.text.replace("Cigna", inputs.customer_name)

    # KPI boxes
    shapes[2].text_frame.paragraphs[0].runs[0].text = fmt_currency_short(calcs.total_3yr_benefits)
    shapes[3].text_frame.paragraphs[0].runs[0].text = calcs.payback_period
    shapes[4].text_frame.paragraphs[0].runs[0].text = fmt_currency_short(calcs.npv)
    shapes[8].text_frame.paragraphs[0].runs[0].text = fmt_roi(calcs.roi_pct)

    # Title — "Executive summary prepared for Cigna"
    shapes[7].text_frame.paragraphs[0].runs[1].text = inputs.customer_name

    # Key Benefits bar chart (inside Group 24, shape[13]).
    # Categories match the benefit prose order above; keep aligned if that changes.
    chart = list(shapes[13].shapes)[0].chart
    data = CategoryChartData()
    data.categories = [
        "DBA Productivity & Labor Cost Avoidance",
        "Faster Incident Resolution (MTTR Reduction)",
        "Avoided Severe Incidents & Downtime Risk",
        "Tool Consolidation & Operational Efficiency",
    ]
    data.add_series("Benefits", (
        calcs.dba_total,
        calcs.mttr_total,
        calcs.sev1_total,
        calcs.tool_total,
    ))
    chart.replace_data(data)

    # Template's data-label format rounds millions to integers ("$3M" for $2.62M).
    # Swap in a 1-decimal format that matches the value axis.
    _DLBL_FMT = '[>999999]"$"0.0,,"M";[>999]"$"0,"K";"$"0'
    for num_fmt in chart._chartSpace.iter(qn("c:numFmt")):
        parent_tag = etree.QName(num_fmt.getparent()).localname
        if parent_tag == "dLbls":
            num_fmt.set("formatCode", _DLBL_FMT)


# ---------------------------------------------------------------------------
# Slide 8 (index 7): Benefits Overview Table
# ---------------------------------------------------------------------------
def update_slide_08(slide, inputs: BvaRequest, calcs: BvaCalculations):
    # Title: "Genius Hub empowers <Customer> to save $X.XM over three years"
    title_runs = slide.shapes[0].text_frame.paragraphs[0].runs
    title_runs[0].text = title_runs[0].text.replace("Cigna", inputs.customer_name)
    title_runs[1].text = fmt_currency_short(calcs.total_3yr_benefits)

    # Table (6 rows × 5 cols). Row order in template:
    #   1: SEV-1, 2: DBA, 3: MTTR, 4: Tool, 5: Total
    table = slide.shapes[1].table

    _set_cell(table, 1, 1, fmt_currency_full(calcs.sev1_yr1))
    _set_cell(table, 1, 2, fmt_currency_full(calcs.sev1_yr2))
    _set_cell(table, 1, 3, fmt_currency_full(calcs.sev1_yr3))
    _set_cell(table, 1, 4, fmt_currency_full(calcs.sev1_total))

    _set_cell(table, 2, 1, fmt_currency_full(calcs.dba_yr1))
    _set_cell(table, 2, 2, fmt_currency_full(calcs.dba_yr2))
    _set_cell(table, 2, 3, fmt_currency_full(calcs.dba_yr3))
    _set_cell(table, 2, 4, fmt_currency_full(calcs.dba_total))

    _set_cell(table, 3, 1, fmt_currency_full(calcs.mttr_yr1))
    _set_cell(table, 3, 2, fmt_currency_full(calcs.mttr_yr2))
    _set_cell(table, 3, 3, fmt_currency_full(calcs.mttr_yr3))
    _set_cell(table, 3, 4, fmt_currency_full(calcs.mttr_total))

    _set_cell(table, 4, 1, fmt_currency_full(calcs.tool_yr1))
    _set_cell(table, 4, 2, fmt_currency_full(calcs.tool_yr2))
    _set_cell(table, 4, 3, fmt_currency_full(calcs.tool_yr3))
    _set_cell(table, 4, 4, fmt_currency_full(calcs.tool_total))

    _set_cell(table, 5, 1, fmt_currency_full(calcs.benefit_yr1))
    _set_cell(table, 5, 2, fmt_currency_full(calcs.benefit_yr2))
    _set_cell(table, 5, 3, fmt_currency_full(calcs.benefit_yr3))
    _set_cell(table, 5, 4, fmt_currency_full(calcs.total_3yr_benefits))

    # Four doughnut charts (shapes 6-9, DBA/MTTR/SEV-1/Tool order).
    # Each shows the pillar's share of 3-year total; center label shows the dollar value.
    total = calcs.total_3yr_benefits
    pillars = [
        (6, calcs.dba_total),
        (7, calcs.mttr_total),
        (8, calcs.sev1_total),
        (9, calcs.tool_total),
    ]
    for shape_idx, pillar_total in pillars:
        outer = slide.shapes[shape_idx]
        inner = list(outer.shapes)[1]
        label_box = list(inner.shapes)[0]
        chart = list(inner.shapes)[1].chart

        share = pillar_total / total if total else 0.0
        data = CategoryChartData()
        data.categories = ["Label 1", "Label 2"]
        data.add_series("Sales", (share, 1 - share))
        chart.replace_data(data)

        label_box.text_frame.paragraphs[0].runs[0].text = fmt_currency_short(pillar_total)


# ---------------------------------------------------------------------------
# Slide 9 (index 8): Why Now?
# ---------------------------------------------------------------------------
def update_slide_09(slide, inputs: BvaRequest, calcs: BvaCalculations):
    tf = slide.shapes[8].text_frame
    tf.paragraphs[0].runs[1].text = inputs.customer_name
    tf.paragraphs[1].runs[0].text = fmt_currency_short(calcs.cost_of_delay_3mo)


# ---------------------------------------------------------------------------
# Slide 12 (index 11): Key Assumptions
# ---------------------------------------------------------------------------

# Row index → short citation for the Source column
_ASSUMPTION_SOURCES = {
    0: "Source",
    1: "Forrester TEI",
    2: "Customer-provided",
    3: "BLS OEWS May 2024",
    4: "IDC Oracle DB Study",
    5: "Gartner IOCS 2024",
    6: "Customer-provided *",
    7: "ITIC 2022; Ponemon 2016",
    8: "Gartner 2024; Forrester TEI",
    9: "New Relic 2025 (n=1,700)",
    10: "Forrester TEI; Capterra",
    11: "Customer-provided *",
    12: "IBM/Ponemon 2023",
    13: "Forrester TEI (IBM Instana)",
}


def update_slide_12(slide, inputs: BvaRequest):
    table = slide.shapes[4].table

    ramp_str = f"{fmt_percentage(RAMP_YR1)}, {fmt_percentage(RAMP_YR2)}, {fmt_percentage(RAMP_YR3)}"
    _set_cell(table, 1, 1, ramp_str)
    _set_cell(table, 2, 1, str(inputs.num_dbas))
    _set_cell(table, 3, 1, fmt_currency_full(inputs.dba_annual_pay))
    _set_cell(table, 4, 1, fmt_percentage(inputs.dba_ops_pct))
    _set_cell(table, 5, 1, fmt_percentage(inputs.ops_reduction_pct))
    _set_cell(table, 6, 1, str(inputs.incidents_per_year))
    _set_cell(table, 7, 1, fmt_currency_full(inputs.cost_per_incident))
    _set_cell(table, 8, 1, fmt_percentage(inputs.mttr_reduction_pct))
    _set_cell(table, 9, 1, str(inputs.num_tools))
    _set_cell(table, 10, 1, fmt_currency_full(inputs.cost_per_tool))
    _set_cell(table, 11, 1, str(inputs.sev1_per_year))
    _set_cell(table, 12, 1, fmt_currency_full(inputs.cost_per_sev1))
    _set_cell(table, 13, 1, fmt_percentage(inputs.sev1_reduction_pct))

    for row_idx, source_text in _ASSUMPTION_SOURCES.items():
        _set_source_cell(table, row_idx, 2, source_text, header=(row_idx == 0))


# ---------------------------------------------------------------------------
# Slide 13 (index 12): DBA Productivity Detail
# ---------------------------------------------------------------------------
def update_slide_13(slide, inputs: BvaRequest, calcs: BvaCalculations):
    slide.shapes[2].text_frame.paragraphs[0].runs[0].text = fmt_currency_short(calcs.dba_total)

    table = slide.shapes[5].table
    _set_cell(table, 0, 2, str(inputs.num_dbas))
    _set_cell(table, 1, 2, fmt_currency_full(inputs.dba_annual_pay))
    _set_cell(table, 2, 2, fmt_percentage(inputs.dba_ops_pct))
    _set_cell(table, 3, 2, fmt_percentage(inputs.ops_reduction_pct))
    _set_cell(table, 4, 2, fmt_currency_full(calcs.total_dba_labor))
    _set_cell(table, 5, 2, fmt_currency_full(calcs.ops_portion))
    _set_cell(table, 7, 2, fmt_percentage(RAMP_YR1))
    _set_cell(table, 7, 3, fmt_percentage(RAMP_YR2))
    _set_cell(table, 7, 4, fmt_percentage(RAMP_YR3))
    _set_cell(table, 8, 2, fmt_currency_full(calcs.dba_yr1))
    _set_cell(table, 8, 3, fmt_currency_full(calcs.dba_yr2))
    _set_cell(table, 8, 4, fmt_currency_full(calcs.dba_yr3))


# ---------------------------------------------------------------------------
# Slide 14 (index 13): MTTR Reduction Detail
# ---------------------------------------------------------------------------
def update_slide_14(slide, inputs: BvaRequest, calcs: BvaCalculations):
    slide.shapes[2].text_frame.paragraphs[0].runs[0].text = fmt_currency_short(calcs.mttr_total)

    table = slide.shapes[5].table
    _set_cell(table, 0, 2, str(inputs.incidents_per_year))
    _set_cell(table, 1, 2, fmt_currency_full(inputs.cost_per_incident))
    _set_cell(table, 2, 2, fmt_percentage(inputs.mttr_reduction_pct))
    _set_cell(table, 3, 2, fmt_currency_full(calcs.annual_incident_cost))
    _set_cell(table, 5, 2, fmt_percentage(RAMP_YR1))
    _set_cell(table, 5, 3, fmt_percentage(RAMP_YR2))
    _set_cell(table, 5, 4, fmt_percentage(RAMP_YR3))
    _set_cell(table, 6, 2, fmt_currency_full(calcs.mttr_yr1))
    _set_cell(table, 6, 3, fmt_currency_full(calcs.mttr_yr2))
    _set_cell(table, 6, 4, fmt_currency_full(calcs.mttr_yr3))


# ---------------------------------------------------------------------------
# Slide 15 (index 14): SEV-1 Avoidance Detail
# ---------------------------------------------------------------------------
def update_slide_15(slide, inputs: BvaRequest, calcs: BvaCalculations):
    slide.shapes[2].text_frame.paragraphs[0].runs[0].text = fmt_currency_short(calcs.sev1_total)

    table = slide.shapes[5].table
    _set_cell(table, 0, 2, str(inputs.sev1_per_year))
    _set_cell(table, 1, 2, fmt_currency_full(inputs.cost_per_sev1))
    _set_cell(table, 2, 2, fmt_percentage(inputs.sev1_reduction_pct))
    _set_cell(table, 3, 2, fmt_currency_full(calcs.baseline_sev1_cost))
    _set_cell(table, 5, 2, fmt_percentage(RAMP_YR1))
    _set_cell(table, 5, 3, fmt_percentage(RAMP_YR2))
    _set_cell(table, 5, 4, fmt_percentage(RAMP_YR3))
    _set_cell(table, 6, 2, fmt_currency_full(calcs.sev1_yr1))
    _set_cell(table, 6, 3, fmt_currency_full(calcs.sev1_yr2))
    _set_cell(table, 6, 4, fmt_currency_full(calcs.sev1_yr3))


# ---------------------------------------------------------------------------
# Slide 16 (index 15): Tool Consolidation Detail
# ---------------------------------------------------------------------------
def update_slide_16(slide, inputs: BvaRequest, calcs: BvaCalculations):
    slide.shapes[2].text_frame.paragraphs[0].runs[0].text = fmt_currency_short(calcs.tool_total)

    table = slide.shapes[5].table
    _set_cell(table, 0, 2, str(inputs.num_tools))
    _set_cell(table, 1, 2, fmt_currency_full(inputs.cost_per_tool))
    _set_cell(table, 3, 2, fmt_percentage(RAMP_YR1))
    _set_cell(table, 3, 3, fmt_percentage(RAMP_YR2))
    _set_cell(table, 3, 4, fmt_percentage(RAMP_YR3))
    _set_cell(table, 4, 2, fmt_currency_full(calcs.tool_yr1))
    _set_cell(table, 4, 3, fmt_currency_full(calcs.tool_yr2))
    _set_cell(table, 4, 4, fmt_currency_full(calcs.tool_yr3))
