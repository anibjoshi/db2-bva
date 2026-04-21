import io
import textwrap

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.util import Emu

from formatting import fmt_currency_short
from models import BvaCalculations

# IBM Design Language colors
IBM_BLUE_HEX = "#0f62fe"
IBM_GRAY_HEX = "#c6c6c6"
IBM_BLUE_LINE_HEX = "#0043ce"

FONT_FAMILY = "Arial"
CHART_DPI = 200


def _find_shape_by_name(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _replace_image(slide, shape_name, img_bytes, match_height=False):
    """Replace an existing image shape's content, preserving position and width.
    If match_height=False (default), height scales from native aspect ratio.
    If match_height=True, both width and height are locked to the original shape."""
    shape = _find_shape_by_name(slide, shape_name)
    if shape is None:
        return
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    el = shape._element
    el.getparent().remove(el)
    if match_height:
        slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width=width, height=height)
    else:
        slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width=width)


# ---------------------------------------------------------------------------
# Slide 2: taxonomy_chart — donut infographic with 4 categories
# ---------------------------------------------------------------------------
def create_slide2_chart(slide, calcs: BvaCalculations):
    categories = [
        ("DBA Productivity\n& Labor Cost\nAvoidance", calcs.dba_total),
        ("Faster Incident\nResolution (MTTR\nReduction)", calcs.mttr_total),
        ("Avoided Severe\nIncidents &\nDowntime Risk", calcs.sev1_total),
        ("Tool\nConsolidation &\nOperational Ef...", calcs.tool_total),
    ]

    total = calcs.total_3yr_benefits
    fig, axes = plt.subplots(1, 4, figsize=(14, 4.5))
    fig.patch.set_facecolor("white")

    for i, (label, value) in enumerate(categories):
        ax = axes[i]
        proportion = value / total if total > 0 else 0
        # Donut chart: filled portion + remainder
        sizes = [proportion, 1 - proportion]
        colors_donut = [IBM_BLUE_HEX, "#e0e0e0"]
        wedges, _ = ax.pie(
            sizes, startangle=90, counterclock=False,
            colors=colors_donut,
            wedgeprops=dict(width=0.3, edgecolor="white"),
        )
        # Center label
        ax.text(0, 0, fmt_currency_short(value),
                ha="center", va="center",
                fontsize=13, fontfamily=FONT_FAMILY, fontweight="normal")
        # Category label below
        ax.set_title(label, fontsize=16, fontfamily=FONT_FAMILY,
                     fontweight="normal", pad=8, y=-0.15,
                     linespacing=1.3)
        ax.set_aspect("equal")

    plt.tight_layout(rect=[0, 0.02, 1, 0.92])

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=CHART_DPI, bbox_inches="tight",
                facecolor="white", edgecolor="none")
    plt.close(fig)
    buf.seek(0)

    _replace_image(slide, "taxonomy_chart", buf.read())


# ---------------------------------------------------------------------------
# Slide 3: summary_chart — bar + cumulative line
# ---------------------------------------------------------------------------
def create_slide3_chart(slide, calcs: BvaCalculations):
    years = ["Year 1", "Year 2", "Year 3"]
    benefits = [calcs.benefit_yr1, calcs.benefit_yr2, calcs.benefit_yr3]
    cumulative = [
        calcs.benefit_yr1,
        calcs.benefit_yr1 + calcs.benefit_yr2,
        calcs.benefit_yr1 + calcs.benefit_yr2 + calcs.benefit_yr3,
    ]

    fig, ax = plt.subplots(figsize=(10, 6))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    x = np.arange(len(years))
    bar_width = 0.45

    # Bars
    bars = ax.bar(x, benefits, bar_width, color=IBM_BLUE_HEX, label="Benefit", zorder=3)

    # Cumulative line
    ax.plot(x, cumulative, color=IBM_BLUE_LINE_HEX, marker="o",
            markersize=8, markerfacecolor="white", markeredgecolor=IBM_BLUE_LINE_HEX,
            markeredgewidth=2, linewidth=2, label="Cumulative Benefit", zorder=4)

    # Data labels on line
    for i, val in enumerate(cumulative):
        ax.annotate(fmt_currency_short(val),
                    (x[i], val), textcoords="offset points",
                    xytext=(0, 12), ha="center", fontsize=10,
                    fontfamily=FONT_FAMILY)

    # Title
    ax.set_title(f"3-Year Total = {fmt_currency_short(calcs.total_3yr_benefits)}",
                 fontsize=14, fontfamily=FONT_FAMILY, fontweight="normal", pad=15)

    ax.set_xticks(x)
    ax.set_xticklabels(years, fontsize=11, fontfamily=FONT_FAMILY)

    # Y-axis formatting
    max_val = max(cumulative) * 1.25
    ax.set_ylim(0, max_val)
    ax.yaxis.set_major_formatter(
        mticker.FuncFormatter(lambda v, _: fmt_currency_short(v) if v > 0 else "$0")
    )
    ax.tick_params(axis="y", labelsize=9)

    # Grid
    ax.grid(axis="y", linestyle="-", alpha=0.3, zorder=0)
    ax.set_axisbelow(True)

    # Legend
    ax.legend(loc="lower center", bbox_to_anchor=(0.5, -0.15),
              ncol=2, fontsize=10, frameon=False,
              prop={"family": FONT_FAMILY})

    # Clean up spines
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    plt.tight_layout()

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=CHART_DPI, bbox_inches="tight",
                facecolor="white", edgecolor="none")
    plt.close(fig)
    buf.seek(0)

    _replace_image(slide, "summary_chart", buf.read(), match_height=True)


# ---------------------------------------------------------------------------
# Slide 6: Update existing BAR_CLUSTERED chart in Group 24
# ---------------------------------------------------------------------------
def update_slide6_chart(slide, calcs: BvaCalculations):
    group = slide.shapes[16]  # Group 24
    chart_shape = group.shapes[0]  # top_benefits_chart_std

    chart_data = CategoryChartData()
    chart_data.categories = [
        "DBA Productivity & Labor Cost Avoidance",
        "Faster Incident Resolution (MTTR Reduction)",
        "Avoided Severe Incidents & Downtime Risk",
        "Tool Consolidation & Operational Efficiency",
    ]
    chart_data.add_series("Benefits", (
        calcs.dba_total,
        calcs.mttr_total,
        calcs.sev1_total,
        calcs.tool_total,
    ))
    chart_shape.chart.replace_data(chart_data)
